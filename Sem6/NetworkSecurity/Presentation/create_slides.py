from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Create a presentation object
prs = Presentation()

# Define colors
# Dark Blue Background
DARK_BLUE = RGBColor(10, 25, 47) 
# Neon Green/Cyan for Titles
NEON_GREEN = RGBColor(100, 255, 218)
# White for Text
WHITE = RGBColor(255, 255, 255)
# Light Grey for Subtitles
LIGHT_GREY = RGBColor(200, 200, 200)

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    # Set title color
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = NEON_GREEN
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
    
    subtitle.text = subtitle_text
    # Set subtitle color
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = WHITE
        paragraph.font.name = 'Arial'

def add_content_slide(prs, title_text, bullet_points, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = NEON_GREEN
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
    
    # Body Content
    body_shape = slide.placeholders[1]
    
    # If image is present, adjust text box width FIRST
    if image_path and os.path.exists(image_path):
        body_shape.width = Inches(5.0)
        body_shape.left = Inches(0.5) # Explicitly set left margin
        body_shape.top = Inches(1.5)  # Align with typical title bottom
        
        # Add Image
        left = Inches(5.8) # Slight gap from text
        top = Inches(2.0)
        height = Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, height=height)

    tf = body_shape.text_frame
    tf.text = "" # clear existing text (leaves 1 empty paragraph)

    # Use the first paragraph if it exists, otherwise add one
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    # Add points
    for i, point in enumerate(bullet_points):
        # For the first point, use the 'p' we just grabbed/created
        # For subsequent points, add new paragraph
        if i > 0:
            p = tf.add_paragraph()
        
        p.text = point
        p.font.color.rgb = WHITE
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.space_after = Pt(14)
        p.level = 0

# --- Content ---

# Slide 1: Title
add_title_slide(
    prs,
    "Evolution of Endpoint Detection and Response (EDR)",
    "A Comprehensive Review\n\nAuthors: Harpreet Kaur et al."
)

# Slide 2: Introduction
add_content_slide(
    prs,
    "Introduction to EDR",
    [
        "Definition: Security solution monitoring end-user devices to detect and respond to threats.",
        "Core Function: Continuous telemetry collection (processes, network, connections).",
        "Proactive Nature: Unlike reactive Antivirus, EDR hunts for suspicious behavior.",
        "Importance: Critical defense against sophisticated, persistent cyber-attacks."
    ]
)

# Slide 3: Evolution Timeline
timeline_img = "/home/praneesh/.gemini/antigravity/brain/53a5b90c-0956-447e-b582-b6683bb42bdc/edr_evolution_timeline_1771267909702.png"
add_content_slide(
    prs,
    "The Evolution Timeline (2005-2023)",
    [
        "2005-2010 (Foundation): Signatures, Firewalls, Manual Response.",
        "2011-2015 (Behavioral): Logic-based detection, Sandboxing, IoCs.",
        "2016-2020 (Intelligence): Machine Learning, UEBA, File-less detection.",
        "2021-2023 (Advanced): XDR, AI-powered, Zero Trust, Automated Response."
    ],
    image_path=timeline_img
)

# Slide 4: ML & AI
ai_img = "/home/praneesh/.gemini/antigravity/brain/53a5b90c-0956-447e-b582-b6683bb42bdc/ai_in_cybersecurity_1771267934708.png"
add_content_slide(
    prs,
    "Machine Learning & AI in EDR",
    [
        "Role: Analyze massive data volumes impossible for humans.",
        "Pattern Recognition: Establish baselines to flag anomalies.",
        "Zero-Day Detection: Identify unknown threats via intent/behavior.",
        "Challenges: False positives, adversarial attacks, model drift."
    ],
    image_path=ai_img
)

# Slide 5: Advanced Concepts (DOPS)
add_content_slide(
    prs,
    "Advanced Concepts: DOPS",
    [
        "DOPS (Deep Ocean Protection System): Enhances EDR capabilities.",
        "Architecture: Endpoint Service Pack collects data; Microservice processes it.",
        "Innovation: Image-Based Malware Detection.",
        "How it works: Converts binary code to images -> logic detects visual patterns of malware -> faster/more accurate."
    ]
)

# Slide 6: XDR
xdr_img = "/home/praneesh/.gemini/antigravity/brain/53a5b90c-0956-447e-b582-b6683bb42bdc/xdr_architecture_1771267977317.png"
add_content_slide(
    prs,
    "Extended Detection and Response (XDR)",
    [
        "The Future: Evolving beyond endpoints.",
        "Unified Visibility: Network + Cloud + Email + Identity + Endpoint.",
        "Holistic Defense: Correlates signals across layers.",
        "Automated Response: Remediate threats across the entire ecosystem."
    ],
    image_path=xdr_img
)

# Slide 7: Conclusion
add_content_slide(
    prs,
    "Conclusion",
    [
        "Summary: EDR has shifted from signatures to AI-driven behavioral analysis.",
        "Key Driver: Validating AI/ML is essential for modern defense.",
        "Outlook: Integration (XDR) and Automation are the future.",
        "Takeaway: As threats evolve, our defenses must evolve faster."
    ]
)

# Save
output_path = "Evolution_of_EDR_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {os.path.abspath(output_path)}")
